import base64
import logging
from typing import List
import pptx
from pptx.util import Inches, Pt
import os
from settings.settings import settings
from server.utils.errors import INTERNAL_SERVER_ERROR_HTTPEXCEPTION

logger = logging.getLogger(__name__)
# Define custom formatting options
TITLE_FONT_SIZE = Pt(30)
SLIDE_FONT_SIZE = Pt(16)

base_url = settings().server.base_url
async def create_presentation(topic:str, slide_titles:List[str], slide_contents:List[str])->str:
    try:
        prs = pptx.Presentation()
        slide_layout = prs.slide_layouts[1]

        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = topic

        for slide_title, slide_content in zip(slide_titles, slide_contents):
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = slide_title
            slide.shapes.placeholders[1].text = slide_content

            # Customize font size for titles and content
            slide.shapes.title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        paragraph.font.size = SLIDE_FONT_SIZE
        save_path = os.path.join(os.getcwd(), f"static/{topic}_presentation.pptx")
        prs.save(save_path)
        public_path = f"{base_url}/static/{topic}_presentation.pptx"
        return public_path
    except Exception as e:
        logger.error(e)
        raise INTERNAL_SERVER_ERROR_HTTPEXCEPTION(str(e))